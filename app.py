from dotenv import load_dotenv
load_dotenv()

from flask import Flask, request, render_template, send_file, jsonify
from io import BytesIO
from PyPDF2 import PdfReader
from fpdf import FPDF
import google.generativeai as genai
import unicodedata
from google.oauth2 import service_account
from googleapiclient.discovery import build
import os
import re
import time
from functools import wraps
import requests
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
import json

app = Flask(__name__)

# Rate limiting
request_count = {}
RATE_LIMIT = 50  # requests per hour per IP

def rate_limit(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        ip = request.remote_addr
        current_time = time.time()
        
        # Clean old entries
        if ip in request_count:
            request_count[ip] = [t for t in request_count[ip] if current_time - t < 3600]
        else:
            request_count[ip] = []
        
        # Check limit
        if len(request_count[ip]) >= RATE_LIMIT:
            return "Rate limit exceeded. Please try again in an hour.", 429
        
        request_count[ip].append(current_time)
        return f(*args, **kwargs)
    return decorated_function

# Configure APIs
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

# Google Form and Drive API scopes
SCOPES_FORMS = ['https://www.googleapis.com/auth/forms.body']
SCOPES_DRIVE = ['https://www.googleapis.com/auth/drive']


SERVICE_ACCOUNT_JSON = os.getenv("SERVICE_ACCOUNT_JSON")

# Authenticate Google services
try:
    if SERVICE_ACCOUNT_JSON:
        service_info = json.loads(SERVICE_ACCOUNT_JSON)
        creds_forms = service_account.Credentials.from_service_account_info(service_info, scopes=SCOPES_FORMS)
        creds_drive = service_account.Credentials.from_service_account_info(service_info, scopes=SCOPES_DRIVE)
        form_service = build('forms', 'v1', credentials=creds_forms)
        drive_service = build('drive', 'v3', credentials=creds_drive)
        GOOGLE_FORMS_ENABLED = True
    else:
        print("SERVICE_ACCOUNT_JSON not set")
        GOOGLE_FORMS_ENABLED = False
    
except Exception as e:
    print(f"Google Forms integration disabled: {e}")
    GOOGLE_FORMS_ENABLED = False

def validate_email(email):
    """Validate email format"""
    email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    return re.match(email_pattern, email) is not None

def extract_text_from_pdf(file_path):
    """Extract text from PDF"""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text: {e}")
        return ""
def extract_text_from_docx(file_path):
    """Extract text from Word document"""
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return ""

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return ""

def extract_text_from_file(file_path, filename):
    """Extract text based on file extension"""
    ext = filename.lower().split('.')[-1]
    if ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif ext == 'docx':
        return extract_text_from_docx(file_path)
    elif ext == 'pptx':
        return extract_text_from_pptx(file_path)
    elif ext == 'txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        return ""
    
def generate_with_groq(prompt, query):
    """Generate using Groq API (super fast!)"""
    if not GROQ_API_KEY:
        return None
    
    try:
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": "llama-3.3-70b-versatile",  # Best model for quality
                "messages": [
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": query}
                ],
                "temperature": 0.7,
                "max_tokens": 8192
            },
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            return result['choices'][0]['message']['content']
        else:
            print(f"Groq API error: {response.status_code} - {response.text}")
            return None
    except Exception as e:
        print(f"Groq API exception: {e}")
        return None

def generate_with_gemini(prompt, query):
    """Generate using Gemini API"""
    if not GEMINI_API_KEY:
        return None
    
    try:
        models_to_try = [
            "gemini-2.0-flash"
        ]
        
        for model_name in models_to_try:
            try:
                model = genai.GenerativeModel(model_name=model_name)
                full_query = f"{prompt}\n\n{query}"
                
                response = model.generate_content(
                    full_query,
                    generation_config=genai.GenerationConfig(
                        temperature=0.7,
                        top_p=0.95,
                        top_k=40,
                        max_output_tokens=8192,
                    )
                )
                
                print(f"Using Gemini model: {model_name}")
                return response.text.strip()
            except Exception as e:
                print(f"Gemini {model_name} failed: {e}")
                continue
        
        return None
    except Exception as e:
        print(f"Gemini API exception: {e}")
        return None

def generate_questions_with_answers(prompt, question_type, num_questions):
    """Generate questions WITH answer keys using multiple AI providers"""
    
    # Truncate very long prompts
    if len(prompt) > 15000:
        prompt = prompt[:15000] + "..."
    
    system_instruction = """You are an expert educator creating exam questions with answer keys.
    CRITICAL RULES:
    1. NEVER mention "the text", "the passage", "the document" or similar
    2. Write professional, standalone exam questions
    3. ALWAYS provide correct answers
    4. Use clear, academic language"""
    
    if question_type.lower() == "mcq":
        query = f"""Create {num_questions} multiple-choice questions with answer keys.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Question]
   a) [Option]
   b) [Option]
   c) [Option]
   d) [Option]
   **Answer: [letter]**

2. [Question]
   a) [Option]
   b) [Option]
   c) [Option]
   d) [Option]
   **Answer: [letter]**

REQUIREMENTS:
- Direct questions (no "based on text")
- Exactly 4 options per question
- Mark correct answer clearly
- Mix difficulty levels
- Generate EXACTLY {num_questions} questions"""

    elif question_type.lower() == "true/false":
        query = f"""Create {num_questions} true/false questions with answers.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Statement] (True/False)
   **Answer: [True/False]**

2. [Statement] (True/False)
   **Answer: [True/False]**

REQUIREMENTS:
- Clear statements (no source references)
- Mark correct answer
- Balance true and false
- Generate EXACTLY {num_questions} questions"""

    elif question_type.lower() == "fill-ups":
        query = f"""Create {num_questions} fill-in-the-blank questions with answers.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Sentence with ________]
   **Answer: [correct word/phrase]**

2. [Sentence with ________]
   **Answer: [correct word/phrase]**

REQUIREMENTS:
- Use ________ for blank
- No source references
- Provide correct answer
- Generate EXACTLY {num_questions} questions"""

    elif question_type.lower() == "short answer":
        query = f"""Create {num_questions} short answer questions with model answers.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Question]
   **Model Answer: [2-3 sentence answer with key points]**

2. [Question]
   **Model Answer: [2-3 sentence answer with key points]**

REQUIREMENTS:
- Direct questions
- Provide concise model answers
- Generate EXACTLY {num_questions} questions"""

    elif question_type.lower() == "subjective":
        query = f"""Create {num_questions} essay questions with answer guidelines.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Question]
   **Answer Guidelines: [Key points to cover]**

2. [Question]
   **Answer Guidelines: [Key points to cover]**

REQUIREMENTS:
- Analytical questions
- Provide answer guidelines
- Generate EXACTLY {num_questions} questions"""

    else:  # Numerical
        query = f"""Create {num_questions} numerical/problem-solving questions with solutions.

CONTENT:
{prompt}

FORMAT (STRICTLY FOLLOW):
1. [Problem with all data needed]
   **Solution: [Step-by-step solution with final answer]**

2. [Problem with all data needed]
   **Solution: [Step-by-step solution with final answer]**

REQUIREMENTS:
- Complete problem statements
- Provide full solutions
- Generate EXACTLY {num_questions} questions"""

    try:
        # Try multiple models in order of preference
        models_to_try = [
            "gemini-2.0-flash"
        ]
        
        model = None
        last_error = None
        
        for model_name in models_to_try:
            try:
                model = genai.GenerativeModel(model_name=model_name)
                print(f"Using model: {model_name}")
                break
            except Exception as e:
                last_error = e
                continue
        
        if model is None:
            raise Exception(f"No available models. Last error: {last_error}")
        
        # Prepend system instruction to the query
        full_query = f"{system_instruction}\n\n{query}"
        
        response = model.generate_content(
            full_query,
            generation_config=genai.GenerationConfig(
                temperature=0.7,
                top_p=0.95,
                top_k=40,
                max_output_tokens=8192,
            )
        )
        
        questions_text = response.text.strip()
        
        # Clean up
        questions_text = re.sub(r'^(Here are|Here\'s|Below are|Sure).*?[:.\n]', '', questions_text, flags=re.IGNORECASE | re.MULTILINE)
        questions_text = re.sub(r'^#.*$', '', questions_text, flags=re.MULTILINE)
        
        return questions_text.strip()
        
    except Exception as e:
        print(f"Error generating questions: {e}")
        raise

def save_questions_to_pdf(raw_questions, subject, marks, include_answers=True):
    """Save questions to PDF with optional answer key"""
    
    # DEBUG: Print to see what's happening
    print(f"DEBUG: include_answers = {include_answers}, type = {type(include_answers)}")
    
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    # First, clean the raw questions to separate questions from answers
    # Split by question numbers
    question_pattern = r'(\d+\.)'
    parts = re.split(question_pattern, raw_questions)
    
    questions_only = []
    answers_only = []
    
    i = 1
    while i < len(parts):
        if re.match(r'\d+\.', parts[i]):
            q_num = parts[i]
            q_content = parts[i + 1] if i + 1 < len(parts) else ""
            
            # Split content into question part and answer part
            # Look for answer markers
            answer_patterns = [
                r'\*\*Answer:.*',
                r'\*\*Solution:.*',
                r'\*\*Model Answer:.*',
                r'\*\*Answer Guidelines:.*',
                r'Answer:.*$',
                r'Solution:.*$',
            ]
            
            question_text = q_num + q_content
            answer_text = ""
            
            for pattern in answer_patterns:
                match = re.search(pattern, question_text, re.IGNORECASE | re.DOTALL)
                if match:
                    answer_text = match.group(0).replace('**', '').strip()
                    question_text = question_text[:match.start()].strip()
                    break
            
            if question_text.strip():
                questions_only.append(question_text.strip())
                answers_only.append(answer_text.strip() if answer_text else "")
            
            i += 2
        else:
            i += 1
    
    # Page 1: Questions ONLY (no answers ever)
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, "ASSESSMENT PAPER", ln=True, align="C")
    pdf.ln(5)
    
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(0, 8, f"Subject: {subject}", ln=True)
    pdf.cell(0, 8, f"Maximum Marks: {marks}", ln=True)
    pdf.cell(0, 8, f"Time: 3 Hours", ln=True)
    pdf.ln(5)
    
    pdf.set_font("Arial", 'I', 10)
    pdf.multi_cell(0, 6, "Instructions: Read all questions carefully. Write your answers clearly.")
    pdf.ln(5)
    
    pdf.set_draw_color(200, 200, 200)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(5)
    
    pdf.set_font("Arial", size=11)
    for q in questions_only:
        # Extra cleaning - remove ANY remaining answer text
        clean_q = q
        clean_q = re.sub(r'\*\*Answer:.*?\*\*', '', clean_q, flags=re.IGNORECASE | re.DOTALL)
        clean_q = re.sub(r'\*\*Solution:.*?\*\*', '', clean_q, flags=re.IGNORECASE | re.DOTALL)
        clean_q = re.sub(r'\*\*Model Answer:.*?\*\*', '', clean_q, flags=re.IGNORECASE | re.DOTALL)
        clean_q = re.sub(r'\*\*Answer Guidelines:.*?\*\*', '', clean_q, flags=re.IGNORECASE | re.DOTALL)
        # Also catch answers without closing **
        clean_q = re.sub(r'\*\*Answer:.*', '', clean_q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Solution:.*', '', clean_q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Model Answer:.*', '', clean_q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Answer Guidelines:.*', '', clean_q, flags=re.IGNORECASE)
        # Catch plain "Answer:" without asterisks
        clean_q = re.sub(r'^Answer:.*$', '', clean_q, flags=re.IGNORECASE | re.MULTILINE)
        clean_q = re.sub(r'^Solution:.*$', '', clean_q, flags=re.IGNORECASE | re.MULTILINE)
        
        clean_q = clean_q.strip()
        if clean_q:
            pdf.multi_cell(0, 6, remove_special_characters(clean_q))
            pdf.ln(3)
    
    # Page 2: Answer Key - ONLY if include_answers is True
    if include_answers == True:  # Explicit check
        print("DEBUG: Adding answer key page")
        pdf.add_page()
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(0, 10, "ANSWER KEY", ln=True, align="C")
        pdf.ln(10)
        
        pdf.set_font("Arial", size=11)
        for i, answer in enumerate(answers_only, 1):
            if answer.strip():
                pdf.set_font("Arial", 'B', 11)
                pdf.cell(0, 6, f"Question {i}:", ln=True)
                pdf.set_font("Arial", size=10)
                pdf.multi_cell(0, 6, remove_special_characters(answer))
                pdf.ln(3)
    else:
        print("DEBUG: Skipping answer key page")
    
    pdf_bytes = BytesIO()
    pdf_string = pdf.output(dest='S').encode('latin1')
    pdf_bytes.write(pdf_string)
    pdf_bytes.seek(0)
    return pdf_bytes


def save_questions_to_docx(raw_questions, subject, marks, include_answers=True):
    """Save questions to Word document (editable)"""
    doc = Document()
    
    # Title
    title = doc.add_heading('ASSESSMENT PAPER', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Info
    doc.add_paragraph(f"Subject: {subject}")
    doc.add_paragraph(f"Maximum Marks: {marks}")
    doc.add_paragraph(f"Time: 3 Hours")
    doc.add_paragraph()
    
    instructions = doc.add_paragraph("Instructions: Read all questions carefully. Write your answers clearly.")
    instructions.italic = True
    doc.add_paragraph("_" * 70)
    doc.add_paragraph()
    
    # Parse questions and answers
    question_pattern = r'(\d+\.)'
    parts = re.split(question_pattern, raw_questions)
    
    questions_only = []
    answers_only = []
    
    i = 1
    while i < len(parts):
        if re.match(r'\d+\.', parts[i]):
            q_num = parts[i]
            q_content = parts[i + 1] if i + 1 < len(parts) else ""
            question_text = q_num + q_content
            answer_text = ""
            
            for pattern in [r'\*\*Answer:.*', r'\*\*Solution:.*', r'\*\*Model Answer:.*', r'\*\*Answer Guidelines:.*']:
                match = re.search(pattern, question_text, re.IGNORECASE | re.DOTALL)
                if match:
                    answer_text = match.group(0).replace('**', '').strip()
                    question_text = question_text[:match.start()].strip()
                    break
            
            if question_text.strip():
                questions_only.append(question_text.strip())
                answers_only.append(answer_text.strip() if answer_text else "")
            i += 2
        else:
            i += 1
    
    # Add questions
    doc.add_heading('Questions', level=1)
    for q in questions_only:
        clean_q = re.sub(r'\*\*Answer:.*', '', q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Solution:.*', '', clean_q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Model Answer:.*', '', clean_q, flags=re.IGNORECASE)
        clean_q = re.sub(r'\*\*Answer Guidelines:.*', '', clean_q, flags=re.IGNORECASE)
        if clean_q.strip():
            doc.add_paragraph(clean_q.strip())
            doc.add_paragraph()
    
    # Add answers if requested
    if include_answers == True:
        doc.add_page_break()
        answer_title = doc.add_heading('ANSWER KEY', 0)
        answer_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()
        
        for i, answer in enumerate(answers_only, 1):
            if answer.strip():
                ans_para = doc.add_paragraph()
                run = ans_para.add_run(f"Question {i}: ")
                run.bold = True
                ans_para.add_run(answer)
                doc.add_paragraph()
    
    docx_bytes = BytesIO()
    doc.save(docx_bytes)
    docx_bytes.seek(0)
    return docx_bytes

def remove_special_characters(text):
    """Remove special characters for PDF"""
    return unicodedata.normalize('NFKD', text).encode('ascii', 'ignore').decode('utf-8')

def create_google_form(questions, user_email):
    """Create Google Form (questions only, no answers)"""
    if not GOOGLE_FORMS_ENABLED:
        return None
        
    try:
        new_form = {
            "info": {
                "title": "Generated Assessment Form",
                "documentTitle": "Assessment Form"
            }
        }
        form = form_service.forms().create(body=new_form).execute()
        form_id = form['formId']
        form_url = form["responderUri"]

        if user_email:
            share_google_form(form_id, user_email)

        questions_list = parse_questions(questions)
        for q_data in questions_list:
            add_question_to_form(form_service, form_id, q_data['question'], q_data['type'], q_data.get('options'))

        return form_url
    except Exception as e:
        print(f"Error creating form: {e}")
        return None

def parse_questions(questions_text):
    """Parse questions (strip answers for forms)"""
    questions = []
    lines = questions_text.strip().split('\n')
    current_question = None
    current_options = []
    
    for line in lines:
        line = line.strip()
        if not line or '**Answer' in line or '**Solution' in line or '**Model' in line:
            continue
            
        if re.match(r'^\d+\.', line):
            if current_question:
                q_type = "mcq" if current_options else "text"
                questions.append({
                    'question': current_question,
                    'type': q_type,
                    'options': current_options if current_options else None
                })
            
            current_question = re.sub(r'^\d+\.\s*', '', line)
            current_options = []
            
        elif re.match(r'^[a-d]\)', line):
            option_text = re.sub(r'^[a-d]\)\s*', '', line)
            current_options.append(option_text)
    
    if current_question:
        q_type = "mcq" if current_options else "text"
        questions.append({
            'question': current_question,
            'type': q_type,
            'options': current_options if current_options else None
        })
    
    return questions

def share_google_form(form_id, user_email):
    """Share form with user"""
    if not GOOGLE_FORMS_ENABLED:
        return
        
    permission = {
        'type': 'user',
        'role': 'writer',
        'emailAddress': user_email
    }
    try:
        drive_service.permissions().create(fileId=form_id, body=permission, sendNotificationEmail=True).execute()
    except Exception as e:
        print(f"Error sharing form: {e}")

def add_question_to_form(service, form_id, question_text, question_type="text", choices=None):
    """Add question to form"""
    if not question_text.strip():
        return
    
    form = service.forms().get(formId=form_id).execute()
    existing_questions = form.get('items', [])
    question_count = len(existing_questions)

    if question_type == "mcq" and choices:
        question = {
            "requests": [{
                "createItem": {
                    "item": {
                        "title": question_text,
                        "questionItem": {
                            "question": {
                                "required": True,
                                "choiceQuestion": {
                                    "type": "RADIO",
                                    "options": [{"value": choice} for choice in choices]
                                }
                            }
                        }
                    },
                    "location": {"index": question_count}
                }
            }]
        }
    else:
        question = {
            "requests": [{
                "createItem": {
                    "item": {
                        "title": question_text,
                        "questionItem": {
                            "question": {
                                "required": True,
                                "textQuestion": {"paragraph": True}
                            }
                        }
                    },
                    "location": {"index": question_count}
                }
            }]
        }
    
    try:
        service.forms().batchUpdate(formId=form_id, body=question).execute()
    except Exception as e:
        print(f"Error adding question: {e}")

@app.route('/')
def index():
    return render_template('index.html', google_forms_enabled=GOOGLE_FORMS_ENABLED)

@app.route('/generate', methods=['POST'])
@rate_limit
def generate():
    try:
        uploaded_file = request.files.get('pdf_file')
        if not uploaded_file or uploaded_file.filename == '':
            return "No file uploaded", 400

        filename = uploaded_file.filename
        ext = filename.lower().split('.')[-1]
        
        # NEW: Check for supported file types
        if ext not in ['pdf', 'docx', 'pptx', 'txt']:
            return "Unsupported file format. Please upload PDF, DOCX, PPTX, or TXT files.", 400

        question_type = request.form.get('question_type', 'mcq')
        num_questions = int(request.form.get('num_questions', 10))
        output_format = request.form.get('output_format', 'pdf')
        subject = request.form.get('subject', 'General')
        marks = request.form.get('marks', '100')
        user_email = request.form.get('email', '').strip()
        include_answers = request.form.get('include_answers', 'false').lower() == 'true'

        # Validation
        if num_questions < 1 or num_questions > 50:
            return "Number of questions must be between 1 and 50", 400

        if output_format == 'form':
            if not GOOGLE_FORMS_ENABLED:
                return "Google Forms integration is not configured", 400
            if not user_email:
                return "Email address is required for Google Form", 400
            if not validate_email(user_email):
                return "Invalid email format", 400

        # Save temp file
        temp_path = f"temp_{filename}"
        uploaded_file.save(temp_path)

        try:
            # NEW: Use the multi-format extraction function
            file_text = extract_text_from_file(temp_path, filename)
            
            if len(file_text.strip()) < 100:
                return "File content too short or unreadable. Please upload a file with at least 100 characters of text content.", 400

            # Generate with answers
            questions = generate_questions_with_answers(file_text, question_type, num_questions)

            if output_format == 'pdf':
                pdf_output = save_questions_to_pdf(questions, subject, marks, include_answers)
                fname = f'{subject}_Assessment_with_Answers.pdf' if include_answers else f'{subject}_Assessment.pdf'
                return send_file(pdf_output, download_name=fname, as_attachment=True, mimetype='application/pdf')
            
            # NEW: Word document output
            elif output_format == 'docx':
                docx_output = save_questions_to_docx(questions, subject, marks, include_answers)
                fname = f'{subject}_Assessment_with_Answers.docx' if include_answers else f'{subject}_Assessment.docx'
                return send_file(docx_output, download_name=fname, as_attachment=True, 
                               mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            
            elif output_format == 'form':
                try:
                    form_url = create_google_form(questions, user_email)
                    if form_url:
                        return jsonify({'success': True, 'form_url': form_url, 'message': 'Form created successfully!'})
                    else:
                        return jsonify({'success': False, 'message': 'Error creating Google Form. Please check your Google API configuration.'}), 500
                except Exception as e:
                    print(f"Google Form error: {e}")
                    return jsonify({'success': False, 'message': f'Google Form error: {str(e)}'}), 500
            
            else:
                return "Invalid output format", 400
        
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            
    except ValueError as e:
        return f"Invalid input: {str(e)}", 400
    except Exception as e:
        print(f"Error: {str(e)}")
        return f"An error occurred: {str(e)}", 500

@app.route('/health')
def health():
    """Health check endpoint"""
    return jsonify({
        'status': 'healthy',
        'google_forms_enabled': GOOGLE_FORMS_ENABLED,
        'apis_configured': {
            'gemini': bool(GEMINI_API_KEY),
            'groq': bool(GROQ_API_KEY)
        }
    })

if __name__ == '__main__':
    port = int(os.getenv('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)